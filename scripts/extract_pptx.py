"""
pptxファイルからスライドのテキスト・構造情報を抽出するスクリプト。
出力はJSON形式でstdoutに書き出す。

shape_kind フィールド:
  "タイトル"           - タイトルプレースホルダー
  "コンテンツ"         - コンテンツプレースホルダー
  "テキストボックス"   - 独立したテキストボックス
  "図形"               - AutoShape（四角、矢印など）内テキスト
  "表"                 - テーブルセル
  "グラフ"             - チャート（タイトル・軸ラベル）
  "SmartArt"           - SmartArt ダイアグラム内テキスト
  "代替テキスト"       - 図形の Alt テキスト（アクセシビリティ用）
  "その他"             - 上記以外

Usage:
  python extract_pptx.py <path-to-pptx>                 # 全スライド
  python extract_pptx.py <path-to-pptx> --pages 1,3,5  # 指定ページのみ
"""
import sys
import json
import argparse

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
except ImportError:
    print("ERROR: python-pptx is not installed.", file=sys.stderr)
    sys.exit(1)


def parse_pages(pages_str):
    result = set()
    for token in pages_str.split(","):
        token = token.strip()
        if not token:
            continue
        try:
            n = int(token)
            if n < 1:
                print(f"WARNING: ページ番号は1以上で指定してください（無視: {token}）", file=sys.stderr)
            else:
                result.add(n)
        except ValueError:
            print(f"WARNING: 無効なページ番号を無視します: {token}", file=sys.stderr)
    return result


def get_shape_kind(shape):
    """図形の種別を日本語ラベルで返す。レポートの指摘箇所として使用する。"""
    try:
        if getattr(shape, "has_chart", False) and shape.has_chart:
            return "グラフ"
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return "画像"
        if st == MSO_SHAPE_TYPE.TABLE:
            return "表"
        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "テキストボックス"
        if st == MSO_SHAPE_TYPE.PLACEHOLDER:
            # プレースホルダーの種類をさらに細分化
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    return "タイトル"
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    return "サブタイトル"
                else:
                    return "コンテンツ"
            except Exception:
                return "プレースホルダー"
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return "図形"
        if st == MSO_SHAPE_TYPE.GROUP:
            return "グループ図形"
        if st == MSO_SHAPE_TYPE.FREEFORM:
            return "フリーフォーム図形"
        return "その他"
    except Exception:
        return "その他"


def _get_alt_text(shape):
    """図形の Alt テキスト（cNvPr/@descr）を返す。未設定なら None。"""
    try:
        for elem in shape.element.iter():
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            if tag == "cNvPr":
                descr = elem.get("descr", "").strip()
                return descr if descr else None
    except Exception:
        pass
    return None


def extract_text_runs(shape):
    """テキストフレームから段落・ランのテキストを抽出する"""
    paragraphs = []
    if not shape.has_text_frame:
        return paragraphs
    for para in shape.text_frame.paragraphs:
        runs = []
        for run in para.runs:
            runs.append({
                "text": run.text,
                "font_size": run.font.size.pt if run.font.size else None,
                "bold": run.font.bold,
            })
        full_text = "".join(r["text"] for r in runs)
        if full_text.strip():
            paragraphs.append({
                "text": full_text,
                "level": para.level,
                "runs": runs,
            })
    return paragraphs


def iter_shapes(shapes):
    """図形を再帰的に展開する（グループ図形の内部要素を含む）"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


# SmartArt の graphicData URI
_SMARTART_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def is_smartart(shape):
    """図形が SmartArt ダイアグラムかどうかを判定する"""
    try:
        tag = "{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData"
        for elem in shape.element.iter(tag):
            return elem.get("uri") == _SMARTART_URI
    except Exception:
        pass
    return False


def extract_smartart_texts(shape):
    """SmartArt のダイアグラムデータ XML からテキストを抽出する"""
    texts = []
    try:
        dgm_ns = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        rel_ids_elem = None
        for elem in shape.element.iter("{%s}relIds" % dgm_ns):
            rel_ids_elem = elem
            break
        if rel_ids_elem is None:
            return texts

        dm_rid = rel_ids_elem.get("{%s}dm" % r_ns)
        if not dm_rid:
            return texts

        diagram_part = shape.part.related_part(dm_rid)
        root = diagram_part._element

        # <dgm:t> 要素にテキストが格納されている
        for t_elem in root.iter("{%s}t" % dgm_ns):
            if t_elem.text and t_elem.text.strip():
                texts.append(t_elem.text.strip())
    except Exception:
        pass
    return texts


def extract_chart_texts(shape):
    """チャートからタイトル・軸ラベルを (label, text) リストで返す"""
    result = []
    try:
        chart = shape.chart

        # グラフタイトル
        try:
            if chart.has_title and chart.chart_title.has_text_frame:
                text = chart.chart_title.text_frame.text.strip()
                if text:
                    result.append(("グラフタイトル", text))
        except Exception:
            pass

        # カテゴリ軸（X 軸）ラベル
        try:
            ax = chart.category_axis
            if ax.has_title and ax.axis_title.has_text_frame:
                text = ax.axis_title.text_frame.text.strip()
                if text:
                    result.append(("カテゴリ軸ラベル", text))
        except Exception:
            pass

        # 値軸（Y 軸）ラベル
        try:
            ax = chart.value_axis
            if ax.has_title and ax.axis_title.has_text_frame:
                text = ax.axis_title.text_frame.text.strip()
                if text:
                    result.append(("値軸ラベル", text))
        except Exception:
            pass

    except Exception:
        pass
    return result


def extract_slide(slide, slide_number):
    """1枚のスライドから情報を抽出する"""
    title = None
    shapes_data = []
    total_chars = 0

    for shape in iter_shapes(slide.shapes):
        # 画像はスキップ
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue

        kind = get_shape_kind(shape)

        # タイトルを取得
        if kind == "タイトル" and shape.has_text_frame:
            title = shape.text_frame.text.strip()

        # テキストフレームを持つすべての図形（テキストボックス・プレースホルダー・AutoShapeなど）
        if shape.has_text_frame:
            paragraphs = extract_text_runs(shape)
            if paragraphs:
                shape_text = " ".join(p["text"] for p in paragraphs)
                total_chars += len(shape_text)
                shapes_data.append({
                    "shape_name": shape.name,
                    "shape_kind": kind,   # ← 図形種別ラベル
                    "paragraphs": paragraphs,
                })

        # テーブルの処理（セル内テキストを個別に抽出）
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_texts = []
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_text = cell.text_frame.text.strip()
                    if cell_text:
                        table_texts.append(cell_text)
                        total_chars += len(cell_text)
            if table_texts:
                shapes_data.append({
                    "shape_name": shape.name + " (table)",
                    "shape_kind": "表",
                    "paragraphs": [{"text": t, "level": 0, "runs": []} for t in table_texts],
                })

        # グラフ（チャート）の処理：タイトル・軸ラベルを抽出
        if getattr(shape, "has_chart", False) and shape.has_chart:
            chart_items = extract_chart_texts(shape)
            if chart_items:
                paragraphs = [
                    {"text": f"{label}: {text}", "level": 0,
                     "runs": [{"text": f"{label}: {text}", "font_size": None, "bold": None}]}
                    for label, text in chart_items
                ]
                total_chars += sum(len(p["text"]) for p in paragraphs)
                shapes_data.append({
                    "shape_name": shape.name,
                    "shape_kind": "グラフ",
                    "paragraphs": paragraphs,
                })

        # SmartArt ダイアグラムの処理
        if is_smartart(shape):
            smartart_texts = extract_smartart_texts(shape)
            if smartart_texts:
                paragraphs = [
                    {"text": t, "level": 0,
                     "runs": [{"text": t, "font_size": None, "bold": None}]}
                    for t in smartart_texts
                ]
                total_chars += sum(len(p["text"]) for p in paragraphs)
                shapes_data.append({
                    "shape_name": shape.name,
                    "shape_kind": "SmartArt",
                    "paragraphs": paragraphs,
                })

        # 代替テキスト（Alt テキスト）の処理
        # python-pptx は shape.description を公開していないバージョンがあるため
        # XML の cNvPr/@descr 属性を直接参照する
        try:
            alt = _get_alt_text(shape)
            if alt:
                total_chars += len(alt)
                shapes_data.append({
                    "shape_name": shape.name,
                    "shape_kind": "代替テキスト",
                    "paragraphs": [{"text": alt, "level": 0,
                                    "runs": [{"text": alt, "font_size": None, "bold": None}]}],
                })
        except Exception:
            pass

    # ノートペインの抽出
    notes_text = None
    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() or None
    except Exception:
        pass

    return {
        "slide_number": slide_number,
        "title": title or "(タイトルなし)",
        "total_chars": total_chars,
        "shapes": shapes_data,
        "notes": notes_text,
    }


def main():
    parser = argparse.ArgumentParser(
        description="pptxファイルからスライドのテキストと構造情報をJSON形式で抽出する"
    )
    parser.add_argument("pptx_path", help="対象のpptxファイルパス")
    parser.add_argument(
        "--pages",
        default=None,
        help="レビュー対象ページをカンマ区切りで指定（例: 1,3,5）。省略時は全スライドを対象にする",
    )
    args = parser.parse_args()

    target_pages = None
    if args.pages:
        target_pages = parse_pages(args.pages)
        if not target_pages:
            print("ERROR: 有効なページ番号が1つも指定されていません。", file=sys.stderr)
            sys.exit(1)

    try:
        prs = Presentation(args.pptx_path)
    except Exception as e:
        print(f"ERROR: ファイルを開けませんでした: {e}", file=sys.stderr)
        sys.exit(1)

    total_slides = len(prs.slides)

    if target_pages:
        out_of_range = [p for p in sorted(target_pages) if p > total_slides]
        if out_of_range:
            print(
                f"WARNING: 存在しないページ番号を無視します（総スライド数: {total_slides}）: "
                f"{', '.join(str(p) for p in out_of_range)}",
                file=sys.stderr,
            )
            target_pages -= set(out_of_range)
        if not target_pages:
            print("ERROR: 有効なページ番号が1つも残りませんでした。", file=sys.stderr)
            sys.exit(1)

    target_count = len(target_pages) if target_pages else total_slides
    page_label = f"{sorted(target_pages)}" if target_pages else f"全{total_slides}枚"
    print(f"[1/2] テキスト抽出中... (対象: {page_label})", file=sys.stderr)

    slides = []
    extracted_count = 0
    for i, slide in enumerate(prs.slides, start=1):
        if target_pages is None or i in target_pages:
            extracted_count += 1
            print(f"[1/2]   スライド {i}/{total_slides} 処理中 ({extracted_count}/{target_count}枚目)", file=sys.stderr)
            slides.append(extract_slide(slide, i))

    total_chars = sum(s["total_chars"] for s in slides)
    print(f"[1/2] 抽出完了 — {extracted_count}枚, 約{total_chars:,}文字", file=sys.stderr)

    result = {
        "total_slides": total_slides,
        "reviewed_slides": [s["slide_number"] for s in slides],
        "slides": slides,
    }

    import io
    stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
    stdout.write(json.dumps(result, ensure_ascii=False, indent=2))
    stdout.flush()


if __name__ == "__main__":
    main()
