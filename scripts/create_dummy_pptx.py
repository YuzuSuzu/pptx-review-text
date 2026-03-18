"""
テスト用のダミーPowerPointファイルを生成するスクリプト。
意図的に以下の問題を含む資料を作成する：
  - 表記ゆれ（サーバ/サーバー、ユーザ/ユーザー）
  - 文法的ねじれ（主語と述語の不一致）
  - 情報量の偏り（文字数が極端なスライド）
  - 顧客向け専門用語の未説明（RBAC、IaC、APIMなど）
  - トーン不統一（です・ます調とだ・である調の混在）

Usage: python create_dummy_pptx.py <output-path>
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree


def add_slide(prs, layout_index=1):
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def set_title(slide, title_text):
    if slide.shapes.title:
        slide.shapes.title.text = title_text


def add_textbox(slide, text, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(14)
    return txBox


def _add_chart_axis_title(axis_elem, title_text):
    """軸要素に <c:title> を XML で直接追加する（python-pptx は軸タイトルの setter を持たないため）"""
    c_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    title_xml = (
        f'<c:title xmlns:c="{c_ns}" xmlns:a="{a_ns}">'
        f'<c:tx><c:rich>'
        f'<a:bodyPr/><a:lstStyle/>'
        f'<a:p><a:r><a:t>{title_text}</a:t></a:r></a:p>'
        f'</c:rich></c:tx>'
        f'<c:overlay val="0"/>'
        f'</c:title>'
    )
    title_elem = etree.fromstring(title_xml)
    # <c:title> は <c:scaling> の直前に挿入するのが標準的な位置
    scaling_tag = "{%s}scaling" % c_ns
    for i, child in enumerate(axis_elem):
        if child.tag == scaling_tag:
            axis_elem.insert(i, title_elem)
            return
    # scaling が見つからなければ先頭に挿入
    axis_elem.insert(0, title_elem)


def main():
    if len(sys.argv) < 2:
        print("Usage: python create_dummy_pptx.py <output-path>")
        sys.exit(1)

    output_path = sys.argv[1]
    prs = Presentation()

    # ==========================================
    # スライド1: タイトルスライド（問題なし）
    # ==========================================
    slide1 = add_slide(prs, layout_index=0)
    set_title(slide1, "クラウド移行提案書")
    if len(slide1.placeholders) > 1:
        slide1.placeholders[1].text = "株式会社〇〇御中\n2026年3月 システム開発部"

    # ==========================================
    # スライド2: 表記ゆれ（サーバ/サーバー、ユーザ/ユーザー）
    # ==========================================
    slide2 = add_slide(prs, layout_index=1)
    set_title(slide2, "現状のシステム構成")
    content2 = (
        "現在の構成概要\n"
        "・オンプレミスのサーバーは3台稼働しております。\n"
        "・各サーバに接続するユーザー数は最大50名です。\n"
        "・バックアップサーバは週次で取得しています。\n"
        "・管理ユーザはシステム管理者の2名が担当しています。\n"
        "・新サーバー導入により、処理速度が向上する見込みです。\n"
        "・ユーザ認証にはActive Directoryを使用しています。"
    )
    add_textbox(slide2, content2)

    # ==========================================
    # スライド3: 主語・述語のねじれ、因果関係の破綻
    # ==========================================
    slide3 = add_slide(prs, layout_index=1)
    set_title(slide3, "課題と対応方針")
    content3 = (
        "課題\n"
        "・現行システムの課題はパフォーマンスを改善します。\n"
        "  （※意図：課題はパフォーマンスの低下である）\n"
        "・セキュリティリスクのため、コストが増加する。\n"
        "  （※因果関係不明：セキュリティリスクとコスト増加の関連が不明確）\n\n"
        "対応方針\n"
        "・クラウド移行によりコストは増加する見込みです。そのため、移行を推奨します。\n"
        "  （※コスト増加するのに推奨という矛盾）\n"
        "・ネットワーク帯域は拡張を実施する予定です。"
    )
    add_textbox(slide3, content3)

    # ==========================================
    # スライド4: 情報量過多（文字数が非常に多いスライド）
    # ==========================================
    slide4 = add_slide(prs, layout_index=1)
    set_title(slide4, "移行計画の詳細スケジュール")
    content4 = (
        "フェーズ1: 要件定義・現状調査（1ヶ月目〜2ヶ月目）\n"
        "現行システムのインフラ構成、アプリケーション依存関係、ネットワーク構成、セキュリティポリシー、"
        "バックアップ運用手順、ライセンス状況、サポート契約内容を詳細に調査します。また、"
        "ステークホルダーへのヒアリングを通じて移行要件を明確化します。\n\n"
        "フェーズ2: 設計・構築（3ヶ月目〜5ヶ月目）\n"
        "クラウド環境の設計では、ネットワーク設計（VNet、サブネット、NSG）、"
        "IAM設計（ロール定義、ポリシー設定）、ストレージ設計（Blob、ファイル共有）、"
        "監視設計（Azure Monitor、Log Analytics）、バックアップ設計（Recovery Services Vault）を行います。"
        "構築後はIaC（Terraform）を使用したコードレビューを実施します。\n\n"
        "フェーズ3: テスト・検証（6ヶ月目）\n"
        "単体テスト、結合テスト、負荷テスト、セキュリティ診断、UAT（ユーザ受け入れテスト）を順次実施します。"
        "特にRTO（目標復旧時間）とRPO（目標復旧時点）を検証するDRテストを重点的に行います。\n\n"
        "フェーズ4: 移行・切替（7ヶ月目）\n"
        "カットオーバー計画に従い、深夜メンテナンス時間帯に切替を実施します。"
        "ロールバック手順も準備し、問題発生時には即時対応できる体制を整えます。"
    )
    add_textbox(slide4, content4)

    # ==========================================
    # スライド5: 顧客向け専門用語の未説明
    # ==========================================
    slide5 = add_slide(prs, layout_index=1)
    set_title(slide5, "セキュリティ対策方針")
    content5 = (
        "提案するセキュリティ対策\n\n"
        "1. RBACの導入\n"
        "   ユーザーのアクセス権限をロールに基づいて管理します。\n\n"
        "2. APIMによるAPI管理\n"
        "   全てのAPIアクセスをAPIMを経由させ、認証・レート制限を適用します。\n\n"
        "3. IaCによるインフラ管理\n"
        "   TerraformによるIaCでインフラをコードとして管理し、変更履歴を追跡します。\n\n"
        "4. ZTNAの適用\n"
        "   ネットワーク境界に頼らないZTNAモデルを採用し、すべての通信を検証します。\n\n"
        "5. SIEMによる脅威検知\n"
        "   SIEMツールでログを集約し、異常な振る舞いを自動検知します。"
    )
    add_textbox(slide5, content5)

    # ==========================================
    # スライド6: トーン不統一（です・ます ↔ だ・である）
    # ==========================================
    slide6 = add_slide(prs, layout_index=1)
    set_title(slide6, "費用対効果")
    content6 = (
        "コスト比較\n\n"
        "現行オンプレミス環境の維持費は年間約500万円です。\n"
        "一方、クラウド移行後の運用コストは年間約350万円と試算される。\n"
        "初期移行費用は約200万円かかりますが、2年目以降はコスト削減効果が見込まれます。\n\n"
        "ROI（投資対効果）\n"
        "移行コストは初年度で回収できない。しかし、3年間の累計では約450万円のコスト削減が期待できます。\n"
        "また、クラウド移行により運用工数も削減され、担当者の作業時間は月20時間程度削減される見込みだ。\n\n"
        "結論\n"
        "費用対効果の観点から、クラウド移行は有効な投資と判断できます。"
    )
    add_textbox(slide6, content6)

    # ==========================================
    # スライド7: 比較的問題の少ないスライド（参考用）
    # ==========================================
    slide7 = add_slide(prs, layout_index=1)
    set_title(slide7, "サポート体制")
    content7 = (
        "移行後のサポート体制\n\n"
        "・24時間365日の監視サービスを提供します。\n"
        "・障害発生時は1時間以内に初動対応を行います。\n"
        "・月次で運用報告書を提出します。\n"
        "・四半期ごとにシステムレビューを実施します。"
    )
    add_textbox(slide7, content7)

    # ノートペインに問題文を追加（表記ゆれテスト: サーバー）
    notes_slide7 = slide7.notes_slide
    notes_tf7 = notes_slide7.notes_text_frame
    notes_tf7.text = (
        "担当者メモ: 本サーバーの監視ツールはZabbixを使用予定。"
        "サーバ設定の詳細は別途資料を参照のこと。"
    )

    # ==========================================
    # スライド8: 各種AutoShape（四角・丸・ダイヤ）＋表
    # ==========================================
    slide8 = add_slide(prs, layout_index=6)  # 空白レイアウト
    set_title(slide8, "システム構成図（図形テスト）")

    # 説明用テキストボックス
    add_textbox(slide8,
        "以下の図はシステム構成を示します。",
        left=Inches(0.5), top=Inches(1.0), width=Inches(9), height=Inches(0.4))

    # ---- AutoShape: 四角形（表記ゆれあり：サーバー） ----
    box_rect = slide8.shapes.add_shape(
        1,  # RECTANGLE
        Inches(0.3), Inches(1.6), Inches(2.5), Inches(1.0)
    )
    tf = box_rect.text_frame
    tf.word_wrap = True
    tf.text = "Webサーバー\n（表記ゆれテスト）"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(11)
    # Alt テキスト：表記ゆれ（サーバー/サーバ）を意図的に混在させる
    box_rect.element.nvSpPr.cNvPr.set(
        "descr",
        "Webサーバーを示す図形。サーバ障害時の代替処理についての補足説明。"
    )

    # ---- AutoShape: 楕円（丸）（専門用語あり：APIM） ----
    box_oval = slide8.shapes.add_shape(
        9,  # OVAL
        Inches(3.0), Inches(1.6), Inches(2.5), Inches(1.0)
    )
    tf = box_oval.text_frame
    tf.word_wrap = True
    tf.text = "APIMゲートウェイ\n（専門用語テスト）"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(11)
    # Alt テキスト：専門用語未説明（APIM、ZTNA）
    box_oval.element.nvSpPr.cNvPr.set(
        "descr",
        "APIゲートウェイ（APIM）を示す。ZTNAポリシーに基づきアクセス制御を行うコンポーネント。"
    )

    # ---- AutoShape: ダイヤ形（IaC専門用語） ----
    box_diamond = slide8.shapes.add_shape(
        4,  # DIAMOND
        Inches(5.7), Inches(1.6), Inches(2.5), Inches(1.0)
    )
    tf = box_diamond.text_frame
    tf.word_wrap = True
    tf.text = "IaC管理\n（Terraform）"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(11)

    # ---- 表（テーブル）：表記ゆれあり ----
    table_shape = slide8.shapes.add_table(
        3, 3,  # 3行3列
        Inches(0.3), Inches(2.9), Inches(9.0), Inches(1.8)
    )
    tbl = table_shape.table
    headers = ["コンポーネント", "サーバー台数", "担当ユーザー"]
    row0_data = ["Webサーバ", "2台", "管理ユーザ"]
    row1_data = ["DBサーバー", "1台", "管理ユーザー"]  # サーバー・ユーザーは表記ゆれ

    for col_idx, text in enumerate(headers):
        cell = tbl.cell(0, col_idx)
        cell.text = text
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.bold = True

    for col_idx, text in enumerate(row0_data):
        cell = tbl.cell(1, col_idx)
        cell.text = text
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11)

    for col_idx, text in enumerate(row1_data):
        cell = tbl.cell(2, col_idx)
        cell.text = text
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(11)

    # ---- グループ図形（GROUP）：再帰処理テスト用 ----
    # python-pptx はグループ作成の高レベルAPIがないため XML で直接構築する
    spTree = slide8.shapes._spTree
    grpSp_xml = (
        '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        "<p:nvGrpSpPr>"
        '<p:cNvPr id="100" name="Group 100"/>'
        "<p:cNvGrpSpPr/>"
        "<p:nvPr/>"
        "</p:nvGrpSpPr>"
        "<p:grpSpPr>"
        "<a:xfrm>"
        '<a:off x="457200" y="4572000"/>'
        '<a:ext cx="2743200" cy="914400"/>'
        '<a:chOff x="457200" y="4572000"/>'
        '<a:chExt cx="2743200" cy="914400"/>'
        "</a:xfrm>"
        "</p:grpSpPr>"
        # 子図形1: 四角
        "<p:sp>"
        "<p:nvSpPr>"
        '<p:cNvPr id="101" name="GroupRect"/>'
        "<p:cNvSpPr/>"
        "<p:nvPr/>"
        "</p:nvSpPr>"
        "<p:spPr>"
        "<a:xfrm>"
        '<a:off x="457200" y="4572000"/>'
        '<a:ext cx="1143000" cy="914400"/>'
        "</a:xfrm>"
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        "</p:spPr>"
        "<p:txBody>"
        "<a:bodyPr/>"
        "<a:lstStyle/>"
        "<a:p><a:r><a:t>グループ内図形A（ZTNAゲートウェイ）</a:t></a:r></a:p>"
        "</p:txBody>"
        "</p:sp>"
        # 子図形2: 丸
        "<p:sp>"
        "<p:nvSpPr>"
        '<p:cNvPr id="102" name="GroupOval"/>'
        "<p:cNvSpPr/>"
        "<p:nvPr/>"
        "</p:nvSpPr>"
        "<p:spPr>"
        "<a:xfrm>"
        '<a:off x="2057400" y="4572000"/>'
        '<a:ext cx="1143000" cy="914400"/>'
        "</a:xfrm>"
        '<a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        "</p:spPr>"
        "<p:txBody>"
        "<a:bodyPr/>"
        "<a:lstStyle/>"
        "<a:p><a:r><a:t>グループ内図形B（サーバ管理）</a:t></a:r></a:p>"
        "</p:txBody>"
        "</p:sp>"
        "</p:grpSp>"
    )
    grpSp_elem = etree.fromstring(grpSp_xml)
    spTree.append(grpSp_elem)

    # ==========================================
    # スライド9: チャート（グラフタイトル・軸ラベルテスト）
    # ==========================================
    slide9 = add_slide(prs, layout_index=6)  # 空白レイアウト
    set_title(slide9, "コスト分析（グラフテスト）")

    # 棒グラフを追加
    chart_data = ChartData()
    chart_data.categories = ["移行前（オンプレ）", "移行後（クラウド）"]
    chart_data.add_series("年間コスト（万円）", (500, 350))

    chart_shape = slide9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.5),
        chart_data
    )
    chart = chart_shape.chart

    # グラフタイトル（表記ゆれ：サーバー）
    chart.has_title = True
    chart.chart_title.text_frame.text = (
        "クラウド移行前後のコスト比較（サーバー費用含む）"
    )

    # 軸ラベルを XML で直接追加（python-pptx は軸タイトルの setter を持たないため）
    _add_chart_axis_title(chart.value_axis._element, "コスト（万円）")
    _add_chart_axis_title(chart.category_axis._element, "環境（サーバ移行前後）")

    prs.save(output_path)
    print(f"ダミーPPTXファイルを作成しました: {output_path}")
    print(f"総スライド数: {len(prs.slides)}")
    print("含まれる意図的な問題点:")
    print("  スライド2: 表記ゆれ（サーバ/サーバー、ユーザ/ユーザー）")
    print("  スライド3: 主語・述語のねじれ、因果関係の矛盾")
    print("  スライド4: 情報量過多（文字数が非常に多い）")
    print("  スライド5: 専門用語の未説明（RBAC、APIM、IaC、ZTNA、SIEM）")
    print("  スライド6: トーン不統一（です・ます <-> だ・である）")
    print("  スライド7: ノートペインに表記ゆれ（サーバー/サーバ）")
    print("  スライド8: 四角・丸・ダイヤのAutoShape＋表（表記ゆれ・専門用語）＋グループ図形＋Altテキスト")
    print("  スライド9: チャート（グラフタイトル・軸ラベル、表記ゆれ）")


if __name__ == "__main__":
    main()
